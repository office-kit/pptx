#!/usr/bin/env python3
"""Regenerates the minimal-fixture PPTX files in this directory using the
vendored python-pptx (under `references/python-pptx`).

Run from the repo root:

    pip3 install --user lxml Pillow
    python3 test/fixtures/minimal/generate.py

The script writes deterministic-looking files but python-pptx embeds a few
context-dependent timestamps in the core properties; that's fine for our
tests, which use a semantic-equivalence comparator, not byte-equal.

Why python-pptx and not our own writer? Because at this phase of @office-kit/pptx
we have no PresentationML authoring yet — we need *some* real PPTX bytes
to load and assert against. python-pptx is the de-facto reference in OSS
and produces files PowerPoint opens without complaint.
"""

import io
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
REPO_ROOT = os.path.abspath(os.path.join(HERE, "..", "..", ".."))

sys.path.insert(0, os.path.join(REPO_ROOT, "references", "python-pptx", "src"))

from PIL import Image  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402


def write_blank() -> None:
    out = os.path.join(HERE, "blank.pptx")
    Presentation().save(out)
    print(f"  wrote {out} ({os.path.getsize(out)} bytes)")


def write_one_text_slide() -> None:
    out = os.path.join(HERE, "one-text-slide.pptx")
    pres = Presentation()
    layout = pres.slide_layouts[5]  # Title Only
    slide = pres.slides.add_slide(layout)
    title = slide.shapes.title
    if title is not None:
        title.text = "Hello, OOXML"
    pres.save(out)
    print(f"  wrote {out} ({os.path.getsize(out)} bytes)")


def write_two_slides_with_text() -> None:
    out = os.path.join(HERE, "two-slides.pptx")
    pres = Presentation()
    layout = pres.slide_layouts[1]  # Title and Content
    s1 = pres.slides.add_slide(layout)
    s1.shapes.title.text = "Slide 1"
    s1.placeholders[1].text = "Body of slide 1."
    s2 = pres.slides.add_slide(layout)
    s2.shapes.title.text = "Slide 2"
    s2.placeholders[1].text = "Body of slide 2."
    pres.save(out)
    print(f"  wrote {out} ({os.path.getsize(out)} bytes)")


def write_one_image_slide() -> None:
    out = os.path.join(HERE, "one-image-slide.pptx")
    pres = Presentation()
    layout = pres.slide_layouts[5]  # Title Only
    slide = pres.slides.add_slide(layout)
    title = slide.shapes.title
    if title is not None:
        title.text = "Image slide"

    # Solid-red 100x100 PNG via Pillow. Embedded as a byte stream so we don't
    # check a separate image file into the fixtures directory.
    buf = io.BytesIO()
    Image.new("RGB", (100, 100), color=(220, 40, 40)).save(buf, format="PNG")
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(2), Inches(2), Inches(3), Inches(3))

    pres.save(out)
    print(f"  wrote {out} ({os.path.getsize(out)} bytes)")


def main() -> None:
    print("Regenerating minimal fixtures...")
    write_blank()
    write_one_text_slide()
    write_two_slides_with_text()
    write_one_image_slide()
    print("Done.")


if __name__ == "__main__":
    main()
